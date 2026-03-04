from pptx import Presentation

def create_ppt(country, growth, inflation):

    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Global Trade Intelligence"
    slide.placeholders[1].text = "Executive Dashboard Overview"

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Country Analysis"

    slide.placeholders[1].text = f"""
Country: {country}
Export Growth: {growth:.2f}%
Inflation: {inflation:.2f}%

Strategic Recommendation:
Focus on export expansion with macroeconomic monitoring.
"""

    file_name = "Trade_Intelligence_Presentation.pptx"
    prs.save(file_name)

    return file_name